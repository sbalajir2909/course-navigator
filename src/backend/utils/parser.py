"""
utils/parser.py
Parses uploaded files (PDF, PPTX, DOCX) into raw text.
"""
from __future__ import annotations

import io
import subprocess
import tempfile
import os
from pathlib import Path


def parse_file(file_bytes: bytes, filename: str) -> str:
    """
    Parse a file into raw text.

    Supports:
      - PDF  → pypdf
      - PPTX → python-pptx
      - DOCX → pandoc subprocess

    Args:
        file_bytes: Raw bytes of the uploaded file.
        filename:   Original filename (used to determine file type).

    Returns:
        Extracted raw text as a single string.

    Raises:
        ValueError: If the file type is unsupported.
        RuntimeError: If parsing fails.
    """
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        return _parse_pdf(file_bytes)
    elif ext == ".pptx":
        return _parse_pptx(file_bytes)
    elif ext in (".docx", ".doc"):
        return _parse_docx(file_bytes, ext)
    elif ext in (".txt", ".md"):
        # Plain text — decode directly
        return file_bytes.decode("utf-8", errors="replace")
    else:
        raise ValueError(f"Unsupported file type: '{ext}'. Supported: pdf, pptx, docx, txt, md")


# ─────────────────────────────────────────────────────────
# PDF
# ─────────────────────────────────────────────────────────

def _parse_pdf(file_bytes: bytes) -> str:
    """Extract text from a PDF using pypdf."""
    try:
        from pypdf import PdfReader
    except ImportError:
        raise RuntimeError("pypdf is not installed. Run: pip install pypdf")

    reader = PdfReader(io.BytesIO(file_bytes))
    pages: list[str] = []

    for page_num, page in enumerate(reader.pages):
        try:
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"[Page {page_num + 1}]\n{text.strip()}")
        except Exception as exc:
            pages.append(f"[Page {page_num + 1} — extraction error: {exc}]")

    if not pages:
        raise RuntimeError("PDF appears to have no extractable text (possibly scanned/image-based).")

    return "\n\n".join(pages)


# ─────────────────────────────────────────────────────────
# PPTX
# ─────────────────────────────────────────────────────────

def _parse_pptx(file_bytes: bytes) -> str:
    """Extract text from a PowerPoint file using python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx")

    prs = Presentation(io.BytesIO(file_bytes))
    slides: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []

        # Extract text from all shapes (text boxes, titles, content)
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())

            # Also handle tables inside slides
            if shape.has_table:
                for row in shape.table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_cells:
                        slide_texts.append(" | ".join(row_cells))

        if slide_texts:
            slide_content = "\n".join(slide_texts)
            slides.append(f"[Slide {slide_num}]\n{slide_content}")

    if not slides:
        raise RuntimeError("PPTX appears to have no extractable text.")

    return "\n\n".join(slides)


# ─────────────────────────────────────────────────────────
# DOCX
# ─────────────────────────────────────────────────────────

def _parse_docx(file_bytes: bytes, ext: str) -> str:
    """Extract text from a DOCX/DOC file using pandoc subprocess."""
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp_in:
        tmp_in.write(file_bytes)
        tmp_in_path = tmp_in.name

    try:
        result = subprocess.run(
            ["pandoc", "--from", "docx", "--to", "plain", "--wrap=none", tmp_in_path],
            capture_output=True,
            text=True,
            timeout=60,
        )
        if result.returncode != 0:
            raise RuntimeError(
                f"pandoc failed (exit {result.returncode}): {result.stderr.strip()}"
            )
        text = result.stdout.strip()
        if not text:
            raise RuntimeError("DOCX appears to have no extractable text.")
        return text
    except FileNotFoundError:
        raise RuntimeError(
            "pandoc is not installed or not in PATH. Install with: apt-get install pandoc"
        )
    finally:
        os.unlink(tmp_in_path)
